#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 eShopping 專案面試用 PPT 簡報
左側：畫面截圖（需手動添加）
右側：功能描述
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """創建 PPT 簡報"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 標題頁
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "eShopping 電商系統"
    subtitle.text = "面試專案簡報\n\n基於 Spring Framework 的完整電商解決方案"
    
    # 添加姓名（在標題頁底部）
    left = Inches(5)
    top = Inches(6.5)
    width = Inches(4)
    height = Inches(0.5)
    author_box = slide.shapes.add_textbox(left, top, width, height)
    author_frame = author_box.text_frame
    author_frame.text = "[請在此處輸入您的姓名]"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(20)
    author_para.font.color.rgb = RGBColor(255, 255, 255)
    author_para.alignment = PP_ALIGN.CENTER
    
    # 專案概述
    slide = create_content_slide(prs, 
        "專案概述",
        "• 技術架構：Spring Framework 5.3 + Spring Security + Hibernate + Thymeleaf\n"
        "• 資料庫：MySQL 8.0\n"
        "• 開發語言：Java 11\n"
        "• 專案類型：Web 應用程式（WAR）\n"
        "• 功能模組：用戶管理、商品管理、購物車、訂單管理、後台管理\n"
        "• 安全機制：Spring Security 認證授權、BCrypt 密碼加密、角色權限控制")
    
    # 1. 首頁
    slide = create_content_slide(prs,
        "1. 首頁功能",
        "• 展示熱門商品、最新商品、推薦商品\n"
        "• 商品分類導航\n"
        "• 用戶登入狀態顯示\n"
        "• 購物車商品數量顯示\n"
        "• 響應式設計，支援多種裝置瀏覽")
    
    # 2. 用戶認證
    slide = create_content_slide(prs,
        "2. 用戶認證系統",
        "• 用戶註冊：帳號、Email、密碼驗證\n"
        "• 用戶登入：Spring Security 認證機制\n"
        "• 忘記密碼：Email 驗證與密碼重設功能\n"
        "• 密碼加密：BCrypt 加密儲存\n"
        "• 角色管理：ADMIN、USER 角色權限控制")
    
    # 3. 商品展示
    slide = create_content_slide(prs,
        "3. 商品展示功能",
        "• 商品列表：網格/列表視圖展示\n"
        "• 商品搜尋：關鍵字搜尋商品名稱與描述\n"
        "• 分類篩選：依商品類型篩選\n"
        "• 商品詳情：完整商品資訊、圖片、庫存狀態\n"
        "• 價格排序：支援價格升序/降序排序")
    
    # 4. 購物車
    slide = create_content_slide(prs,
        "4. 購物車功能",
        "• 添加商品：從商品列表或詳情頁加入購物車\n"
        "• 數量管理：增加、減少、移除商品\n"
        "• 庫存檢查：即時檢查商品庫存狀態\n"
        "• 總金額計算：自動計算購物車總金額\n"
        "• 購物車持久化：資料庫儲存購物車內容")
    
    # 5. 訂單管理
    slide = create_content_slide(prs,
        "5. 訂單管理系統",
        "• 結帳流程：填寫收件資訊、選擇付款方式\n"
        "• 訂單建立：從購物車轉換為訂單\n"
        "• 訂單查詢：用戶可查看自己的訂單歷史\n"
        "• 訂單詳情：完整的訂單資訊與商品明細\n"
        "• 訂單狀態：PENDING、PROCESSING、SHIPPED、DELIVERED、CANCELLED\n"
        "• 付款狀態：UNPAID、PAID、REFUNDED")
    
    # 6. 商品管理（後台）
    slide = create_content_slide(prs,
        "6. 商品管理（後台）",
        "• 商品 CRUD：新增、編輯、刪除商品\n"
        "• 商品資訊：名稱、描述、價格、庫存、分類\n"
        "• 商品狀態：上架/下架管理\n"
        "• 庫存管理：庫存數量與最低庫存閾值設定\n"
        "• 圖片上傳：支援商品圖片上傳（Google Cloud Storage）")
    
    # 7. 用戶管理（後台）
    slide = create_content_slide(prs,
        "7. 用戶管理（後台）",
        "• 用戶列表：查看所有註冊用戶\n"
        "• 用戶編輯：修改用戶資訊、角色、啟用狀態\n"
        "• 角色管理：ADMIN、USER 角色分配\n"
        "• 帳號啟用/停用：控制用戶帳號狀態\n"
        "• 用戶搜尋：依使用者名稱或 Email 搜尋")
    
    # 8. 客戶管理
    slide = create_content_slide(prs,
        "8. 客戶管理",
        "• 客戶資料：姓名、電話、地址等完整資訊\n"
        "• 客戶編輯：修改個人資料\n"
        "• 客戶列表：管理員可查看所有客戶\n"
        "• 客戶關聯：與 User 帳號一對一關聯\n"
        "• 訂單關聯：客戶與訂單的關聯管理")
    
    # 9. 訂單管理（後台）
    slide = create_content_slide(prs,
        "9. 訂單管理（後台）",
        "• 訂單列表：查看所有訂單\n"
        "• 訂單篩選：依狀態、付款狀態、關鍵字搜尋\n"
        "• 訂單詳情：完整的訂單資訊與收件地址\n"
        "• 狀態更新：管理員可更新訂單狀態\n"
        "• 付款狀態：管理付款狀態\n"
        "• 分頁顯示：支援大量訂單的分頁瀏覽")
    
    # 10. 後台儀表板
    slide = create_content_slide(prs,
        "10. 後台儀表板",
        "• 統計資訊：總用戶數、商品數、訂單數\n"
        "• 訂單統計：待處理訂單、已付款訂單數量\n"
        "• 銷售統計：總銷售額計算\n"
        "• 快速入口：商品管理、訂單管理、用戶管理\n"
        "• 數據視覺化：關鍵指標一目了然")
    
    # 技術特色
    slide = create_content_slide(prs,
        "技術特色",
        "• Spring Security：完整的認證授權機制\n"
        "• Hibernate ORM：物件關聯映射，簡化資料庫操作\n"
        "• Thymeleaf 模板：服務端渲染，SEO 友善\n"
        "• RESTful API：REST 風格的 URL 設計\n"
        "• 事務管理：@Transactional 確保資料一致性\n"
        "• 異常處理：完善的錯誤處理機制\n"
        "• 日誌記錄：SLF4J + Logback 日誌系統")
    
    # 資料庫設計
    slide = create_content_slide(prs,
        "資料庫設計",
        "• User：用戶帳號資訊\n"
        "• Customer：客戶詳細資料\n"
        "• Product：商品資訊\n"
        "• Cart / CartItem：購物車與購物車項目\n"
        "• Order / OrderItem：訂單與訂單項目\n"
        "• OrderAddress：訂單收件地址\n"
        "• 關聯設計：一對一、一對多、多對多關係完整設計")
    
    # 總結
    slide = create_content_slide(prs,
        "專案總結",
        "• 完整的電商系統功能實現\n"
        "• 前後台分離的架構設計\n"
        "• 完善的用戶認證與權限管理\n"
        "• 完整的購物車與訂單流程\n"
        "• 響應式設計，支援多裝置\n"
        "• 程式碼結構清晰，易於維護擴展\n"
        "• 遵循 Spring 最佳實踐")
    
    # 保存檔案
    output_file = "eShopping_專案簡報.pptx"
    prs.save(output_file)
    print(f"簡報已生成：{output_file}")
    return output_file

def create_content_slide(prs, title_text, content_text):
    """創建內容投影片（左側截圖，右側文字）"""
    blank_slide_layout = prs.slide_layouts[6]  # 空白版面
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 標題
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.LEFT
    
    # 左側：截圖區域（佔左側 50%）
    left_img = Inches(0.5)
    top_img = Inches(1.2)
    width_img = Inches(4.5)
    height_img = Inches(5.5)
    
    # 添加文字框作為截圖佔位符
    img_placeholder = slide.shapes.add_textbox(left_img, top_img, width_img, height_img)
    img_frame = img_placeholder.text_frame
    img_frame.text = "[請在此處插入畫面截圖]"
    img_para = img_frame.paragraphs[0]
    img_para.font.size = Pt(16)
    img_para.font.color.rgb = RGBColor(128, 128, 128)
    img_para.alignment = PP_ALIGN.CENTER
    
    # 設定文字框背景為淺灰色
    fill = img_placeholder.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # 右側：功能描述（佔右側 50%）
    left_text = Inches(5.2)
    top_text = Inches(1.2)
    width_text = Inches(4.5)
    height_text = Inches(5.5)
    
    text_box = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content_text
    
    # 設定文字格式
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.line_spacing = 1.2
        paragraph.space_after = Pt(6)
    
    return slide

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("錯誤：請先安裝 python-pptx 套件")
        print("安裝指令：pip install python-pptx")
    except Exception as e:
        print(f"發生錯誤：{e}")

